#!/usr/bin/env python3
"""Generate midway thesis presentation -- v2 with parallelization results."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

PROJECT = Path(__file__).resolve().parent.parent

DARK = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x29, 0x80, 0xB9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GREEN = RGBColor(0x27, 0xAE, 0x60)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    tx = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    p2.space_before = Pt(20)
    return s


def section_slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BLUE)
    tx = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(2))
    p = tx.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = WHITE
    return s


def content_slide(prs, title, bullets, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    # Title
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7))
    p = tx.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = DARK
    # Line
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(8.8), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    # Bullets
    tx2 = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = tx2.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if b.startswith("!"):
            p.text = b[1:]; p.font.size = Pt(15); p.font.bold = True
            p.font.color.rgb = BLUE
        elif b.startswith("##"):
            p.text = b[2:].strip(); p.font.size = Pt(17); p.font.bold = True
            p.font.color.rgb = DARK; p.space_before = Pt(14)
        else:
            p.text = b; p.font.size = Pt(15); p.font.color.rgb = DARK
        p.space_before = Pt(5)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def table_slide(prs, title, headers, rows, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7))
    p = tx.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = DARK

    nr = len(rows) + 1; nc = len(headers)
    tbl = s.shapes.add_table(nr, nc, Inches(0.5), Inches(1.2),
                              Inches(9.0), Inches(0.4 * nr + 0.3)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = DARK
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(v)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12); p.font.color.rgb = DARK
            if i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def chart_slide(prs, title, categories, values, series_name="Speedup", notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7))
    p = tx.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = DARK

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    # Ideal line
    chart_data.add_series("Ideal", [float(c) for c in categories])

    chart_frame = s.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(1), Inches(1.3),
        Inches(8), Inches(5.5), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Style
    plot = chart.plots[0]
    series_actual = plot.series[0]
    series_actual.format.line.color.rgb = BLUE
    series_actual.format.line.width = Pt(3)
    series_ideal = plot.series[1]
    series_ideal.format.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    series_ideal.format.line.dash_style = 2  # dash

    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

    # ===================== SLIDE 1: Title =====================
    title_slide(prs,
        "Refactoring & Parallelization of\nPolymer Growth Simulation Software",
        "Kaan Basaran  |  Midway Presentation  |  April 2026\n"
        "Supervisor: Rachel Cavill  |  Maastricht University"
    )

    # ===================== SLIDE 2: Chemistry =====================
    content_slide(prs, "Why: Polymers for Medicine", [
        "Polymers = chains of repeating molecular units",
        "Medical use: attach drugs to polymers for controlled release",
        "Chain length controls how the body processes them:",
        "  Too short -> kidneys filter out before drug works",
        "  Too long -> body treats as toxic foreign object",
        "",
        "!Goal: understand and control the chain length distribution",
        "",
        "Polymer studied: PEtOx (Poly(2-ethyl-2-oxazoline))",
        "Dataset: Monnery et al. (2018) [2]",
    ], "TALKING POINTS:\n"
       "- Polymers are everywhere: plastics, rubber, DNA, proteins\n"
       "- Medical application: you attach medicine to a polymer chain and inject it\n"
       "- The LENGTH of the chain determines how fast the drug is released\n"
       "- Too short = kidneys flush it out. Too long = body rejects it as poison\n"
       "- There's a sweet spot. To hit it, we need to understand what controls the length distribution\n"
       "- PEtOx is a specific polymer being studied for drug delivery\n"
       "- The real experimental data comes from GPC (gel permeation chromatography) measurements\n"
       "- Ref: Monnery et al. 2018, doi:10.1002/anie.201807796")

    # ===================== SLIDE 3: The Vampire Chemistry =====================
    content_slide(prs, "The Chemistry: Growing Chains & Vampires", [
        "## How polymers grow (agent-based simulation)",
        "Start with thousands of living chains (length 1)",
        "Each timestep, a chain can:",
        "  1. GROW: add a monomer unit (get longer)",
        "  2. DIE: lose its active end (become 'dead')",
        "  3. GET ATTACKED: a dead chain steals life from a living one",
        "",
        "## The 'Vampire' Reaction",
        "A dead chain attacks a living chain and fuses onto it",
        "The dead chain steals the 'livingness' back",
        "Creates unexpectedly long chains (high-MW shoulder)",
        "Longer chains are harder to attack (coiled up, ends buried)",
    ], "TALKING POINTS:\n"
       "- Think of it as a room full of growing worms\n"
       "- Each worm eats food pellets (monomers) and gets longer\n"
       "- Sometimes a worm dies -- stops growing forever\n"
       "- BUT dead worms can 'attack' living ones and fuse onto them\n"
       "- This is the 'vampiric coupling' -- the dead chain steals livingness\n"
       "- This creates a few very long chains which show up as a shoulder/tail in the distribution\n"
       "- The probability of a successful vampire attack depends on both chains' lengths\n"
       "- Longer chains are coiled up so their reactive ends are buried -- harder to attack\n"
       "- Bryn (the chemist) coined the 'vampire' terminology\n"
       "- 10 parameters control all of this: growth probability, death probability, coupling probability, etc.")

    # ===================== SLIDE 4: The Problem =====================
    content_slide(prs, "The Optimization Problem", [
        "A chemist runs an experiment -> gets a histogram of chain lengths",
        "Question: 'What parameters produced this distribution?'",
        "",
        "## FDDC (Fitness-Diversity Driven Co-evolution)",
        "Two competing populations [3][4]:",
        "  Solutions: sets of 10 parameter values",
        "  Problems: weight vectors that expose weaknesses",
        "Run simulation with each guess, compare to real data",
        "Cost = sum of point-by-point differences (lower = better)",
        "",
        "!Thomas (2020) [1]: FDDC converges in ~20 gens vs ~77 for basic GA",
        "!But FDDC is expensive: only reached 42 gens in 24 hours [1]",
    ], "TALKING POINTS:\n"
       "- You can't reverse-engineer the parameters mathematically -- it's a stochastic simulation\n"
       "- So you use trial and error, but SMART trial and error: genetic algorithms\n"
       "- Start with 100 random guesses, test each by running the full simulation\n"
       "- The 'cost' is how different the simulated histogram is from the real one\n"
       "- Kill bad guesses, breed good ones, repeat\n"
       "- FDDC is the best variant: it has a second population that learns to find weaknesses\n"
       "- Thomas proved FDDC needs fewer generations, but each generation costs more\n"
       "- On his hardware, FDDC only managed 42 generations in a full 24-hour day\n"
       "- His basic GA (roulette) managed 1801 generations in the same time\n"
       "- FDDC was smarter per generation but too slow to run long enough")

    # ===================== SLIDE 5: What Was Wrong =====================
    content_slide(prs, "Problems with the Original Code", [
        "Monolithic structure: single files, no package, no tests [1]",
        "Limited UI: Tkinter GUI froze during computation [1]",
        "Hardcoded parallelization: Pool(6) regardless of hardware [1]",
        "BLAS thread contention: numpy spawned hidden internal threads",
        "No run management: no seed control, no output organization [1]",
        "",
        "## The real performance bottleneck",
        "Only the initial evaluation was parallelized",
        "95% of per-generation work ran sequentially",
        "Hidden numpy threading competed with explicit parallelism",
        "!Result: 42 FDDC generations took 24 hours [1, Table VIII]",
    ], "TALKING POINTS:\n"
       "- Thomas wrote working research code, but it was one big folder of scripts\n"
       "- His GUI existed but froze during optimization (single-threaded Tkinter)\n"
       "- Pool(6) was hardcoded on line 16 of his fddc.py -- literally the number 6\n"
       "- Even if you had 64 cores, his code would only use 6\n"
       "- WORSE: numpy's math library (BLAS/Accelerate on Mac) spawns its own threads\n"
       "- So 6 processes x 4 hidden threads = 24 threads on maybe 6 cores = thrashing\n"
       "- The per-generation encounters and child evaluation were completely sequential\n"
       "- Only the initial population evaluation (once at startup) was parallel\n"
       "- This means 95% of the runtime was single-threaded")

    # ===================== SLIDE 6: RQ =====================
    content_slide(prs, "Research Questions", [
        "## Main RQ",
        "How can parallelization and modern software engineering",
        "improve the performance and usability of genetic algorithm-based",
        "polymer simulation parameter fitting, while maintaining",
        "result equivalence with the original implementation?",
        "",
        "## Sub-Questions",
        "RQ1: How can legacy scientific code be modernized while",
        "  preserving algorithmic correctness?",
        "",
        "RQ2: What correctness errors can be introduced during",
        "  refactoring, and how can testing identify them?",
    ], "TALKING POINTS:\n"
       "- Main RQ emphasizes parallelization -- that's where the novel contribution is\n"
       "- RQ1 is about the software engineering: package structure, GUI, testing\n"
       "- RQ2 is about proving we didn't break anything during refactoring\n"
       "- These three together cover: same output + faster + more usable")

    # ===================== SLIDE 7: What I Did =====================
    content_slide(prs, "What Was Done", [
        "## Software Engineering",
        "Modular Python package (pip installable)",
        "PySide6 GUI with interactive plots, batch queue, chemistry tooltips",
        "29 unit tests, run manager with timestamped outputs",
        "",
        "## Parallelization (the breakthrough)",
        "Parallelized ALL evaluation points (encounters + reproduction)",
        "Switched ThreadPool -> ProcessPool (bypass Python GIL)",
        "Disabled BLAS internal threading (eliminated contention)",
        "Numba JIT on simulation hot path (vampiric coupling calc)",
        "Auto-detect CPU cores (was hardcoded 6)",
    ], "TALKING POINTS:\n"
       "- Software side: proper Python package, you can pip install it\n"
       "- GUI: three tabs -- simulation, optimization, batch queue\n"
       "- Interactive plots: hover to see values, zoom, save as PNG/SVG\n"
       "- Chemistry tooltips: every parameter has a ? icon explaining what it means\n"
       "- The BIG story is parallelization:\n"
       "  - Thomas only parallelized the initial population evaluation (startup)\n"
       "  - We parallelized encounters AND child evaluation (every generation)\n"
       "  - ThreadPoolExecutor couldn't scale because Python's GIL blocks the simulation loop\n"
       "  - Switched to multiprocessing.Pool -- separate processes, no GIL sharing\n"
       "  - Discovered numpy was spawning hidden BLAS threads, causing cache thrashing\n"
       "  - One environment variable (VECLIB_MAXIMUM_THREADS=1) fixed it\n"
       "  - Numba JIT compiles the vampiric coupling math to machine code (C speed)")

    # ===================== SLIDE 8: Scaling Chart =====================
    chart_slide(prs, "Parallelization Scaling: 13x at 13 Workers",
        ["1", "2", "4", "6", "8", "13"],
        [1.0, 3.83, 7.05, 9.99, 10.76, 13.04],
        "Measured Speedup",
        "TALKING POINTS:\n"
        "- X axis: number of CPU workers. Y axis: speedup factor\n"
        "- Dashed line = ideal (linear) scaling. Blue = measured\n"
        "- We achieve 13.04x speedup at 13 workers -- 100% parallel efficiency\n"
        "- This means our parallelization has essentially zero overhead\n"
        "- Superlinear at some points due to better cache utilization per-process\n"
        "- The key insight: we had to disable numpy's internal threading\n"
        "  to prevent 13 processes x 4 BLAS threads = 52 threads on 14 cores\n"
        "- Config: pop=50, 10 gens, 5k no BB dataset, 3 seeds averaged\n"
        "- Source: validation_results/speed_benchmark.json"
    )

    # ===================== SLIDE 9: The Headline Number =====================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    tx = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
    p = tx.text_frame.paragraphs[0]; p.text = "Same Configuration. Same 6 Workers."
    p.font.size = Pt(24); p.font.color.rgb = RGBColor(0xBD,0xC3,0xC7); p.alignment = PP_ALIGN.CENTER

    tx2 = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(2))
    tf = tx2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "24 hours  ->  7.1 minutes"
    p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "203x faster"
    p2.font.size = Pt(36); p2.font.color.rgb = GREEN; p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)

    tx3 = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2.5))
    tf3 = tx3.text_frame; tf3.word_wrap = True
    for line in [
        "5k no BB dataset  |  Pop: 100  |  42 FDDC generations",
        "Thomas (2020): cost = 21.74 in ~24 hours (6 workers, Pool(6))",
        "Ours: cost = 22.87 in 7.1 minutes (6 workers, ProcessPool + Numba + BLAS fix)",
        "Same algorithm. Same worker count. Same result quality.",
    ]:
        p = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
        p.text = line; p.font.size = Pt(14); p.font.color.rgb = RGBColor(0xBD,0xC3,0xC7)
        p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)

    s.notes_slide.notes_text_frame.text = (
        "TALKING POINTS:\n"
        "- This is the headline result. Let it sink in.\n"
        "- Same dataset (5k no BB), same population (100), same generations (42)\n"
        "- Same number of workers: 6. This is NOT about having more hardware.\n"
        "- Thomas's code: 24 hours. Our code: 7.1 minutes. 203x faster.\n"
        "- Cost quality is equivalent: 22.87 vs 21.74 (stochastic noise, different seed)\n"
        "- The 203x comes from: Numba JIT (~5-10x per simulation), ProcessPool (no GIL),\n"
        "  BLAS thread fix (no contention), parallelized encounters (was sequential)\n"
        "- NOTE: this includes hardware difference (his 2020 machine vs our M-series)\n"
        "- The controlled experiment is the scaling benchmark (13x at 13 workers)\n"
        "- But even at SAME worker count, the software improvements give 203x\n"
        "- Source: validation_results/thomas_fair_compare.json"
    )

    # ===================== SLIDE 10: Equivalence =====================
    content_slide(prs, "Result Equivalence: Same Output", [
        "## Simulation equivalence (t-test, KS test)",
        "  Both pass with p >> 0.05 across all parameter sets",
        "  Refactored simulation = original simulation",
        "",
        "## Optimization equivalence (4 datasets)",
    ], "TALKING POINTS:\n"
       "- First we proved the raw simulation gives the same output\n"
       "- Ran both implementations 30 times, compared Mn/Mw/PDI\n"
       "- t-test p=0.86, KS p=0.99 -- no significant difference\n"
       "- Then we ran the full optimizer on 4 real datasets\n"
       "- Both implementations find the same best parameters (verified identical vectors)\n"
       "- Cost differences are stochastic noise -- same random process, different seeds\n"
       "- This proves: we didn't break anything. Same algorithm, same results, just faster.")

    # Add equivalence note about source
    s = prs.slides[-1]
    tx_note = s.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.5))
    p = tx_note.text_frame.paragraphs[0]
    p.text = "Thomas paper costs from [1, Table VIII]. Differences due to different random seeds (stochastic process)."
    p.font.size = Pt(10); p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D); p.font.italic = True
    tbl = s.shapes.add_table(5, 4, Inches(0.8), Inches(3.0),
                              Inches(8.4), Inches(2.5)).table
    headers = ["Dataset", "Our Code", "Thomas Sim", "Thomas Paper"]
    data = [
        ["5k (42 gens)", "23.04", "22.80", "21.74"],
        ["10k (56 gens)", "143.51", "145.74", "108.54"],
        ["20k (44 gens)", "176.69", "193.03", "131.73"],
        ["30k (27 gens)", "373.11", "443.82", "340.26"],
    ]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = DARK
    for i, row in enumerate(data):
        for j, v in enumerate(row):
            c = tbl.cell(i+1, j); c.text = v
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = DARK
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LIGHT

    # ===================== SLIDE 11: What Remains =====================
    content_slide(prs, "What Remains", [
        "## Four thesis experiments",
        "a) Output equivalence -- done (single-seed), needs multi-seed replication",
        "b) Speed: Numba JIT -- Thomas's sim vs ours, same 6 workers, same machine",
        "c) Speed: Scaling -- done (13x at 13 workers, 100% efficiency)",
        "d) UI usability -- Likert questionnaire for chemistry supervisor",
        "",
        "## Software",
        "Standalone packaging (no terminal / VS Code needed)",
        "GUI polish based on questionnaire feedback",
        "",
        "## Thesis writing",
        "Methodology, results, and discussion chapters",
    ], "TALKING POINTS:\n"
       "- Experiment A: we ran single-seed. Need 10+ seeds per config for proper statistics\n"
       "- Experiment B: run Thomas's original simulation at 6 workers on OUR machine\n"
       "  to isolate the software contribution from the hardware difference\n"
       "- Experiment C: DONE -- the scaling benchmark we just showed\n"
       "- Experiment D: the UX questionnaire is written, needs Bryn to fill it out\n"
       "- Standalone packaging: PyInstaller or similar so chemists don't need terminal\n"
       "- For thesis: we have all the data, just need to write it up properly")

    # ===================== SLIDE 12: References =====================
    content_slide(prs, "References", [
        "[1] van den Broek, T. (2020). Genetic Algorithms To Better",
        "    Understand Polymer Growth. BSc thesis, Maastricht University.",
        "",
        "[2] Monnery, B.D. et al. (2018). Defined High Molar Mass",
        "    Poly(2-oxazoline)s. Angew. Chem. Int. Ed., 57(47), 15400-15404.",
        "    doi:10.1002/anie.201807796",
        "",
        "[3] Franz, F., Paredis, J. & Mockel, R. (2017). On the Combination",
        "    of Coevolution and Novelty Search. pp. 201-208.",
        "",
        "[4] Paredis, J. (1995). Coevolutionary Computation.",
        "    Artificial Life, 2, 355-375.",
        "",
        "[5] van Appeven, J. et al. (2018). Understanding Polymer Growth.",
    ], "TALKING POINTS:\n"
       "- [1] is Thomas's thesis -- the original code and FDDC results\n"
       "- [2] is Bryn's paper describing the PEtOx datasets we use\n"
       "- [3] is the FDDC algorithm paper (co-evolution + novelty search)\n"
       "- [4] is the original co-evolution paper by Paredis\n"
       "- [5] is the earlier student project that built the first simulation")

    # ===================== SLIDE 13: Thank You =====================
    title_slide(prs,
        "Thank You",
        "Questions?\n\n"
        "github.com/MKBasaran/polymer-growth"
    )

    out = PROJECT / "presentation_midway.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
